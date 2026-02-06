"""
Generate Phase 1 presentation for Data Annotation Platform.
Creates a professional PowerPoint covering architecture, data model,
infrastructure, deployment, monitoring, and confidence-building content.

Run: pip install python-pptx && python generate_phase1_presentation.py
Output: Phase1_Data_Annotation_Platform.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Theme: professional dark accent on light
TITLE_COLOR = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_COLOR = RGBColor(0x16, 0xa3, 0x4a)  # green
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
BULLET_COLOR = RGBColor(0x55, 0x55, 0x55)


def set_title(tf, text, size=28):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR


def add_body(shape, lines, bullet=True):
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = BULLET_COLOR if bullet else BODY_COLOR
        p.space_after = Pt(6)
        if bullet:
            p.level = 0


def add_section_title(tf, text):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(12)


def slide_title(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left, top, width = Inches(0.5), Inches(2), Inches(9)
    box = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = BODY_COLOR
        p2.space_before = Pt(12)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ----- Slide 1: Title -----
    slide = slide_title(
        prs,
        "Data Annotation Platform",
        "Phase 1 — Architecture, Data Model & Delivery Confidence",
    )
    box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Technical proposal for stakeholder alignment"
    p.font.size = Pt(14)
    p.font.color.rgb = BULLET_COLOR

    # ----- Slide 2: Phase 1 at a Glance -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Phase 1 — What We Deliver (8–10 weeks)")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    add_body(
        box2,
        [
            "End-to-end flow: Workspace → Project → Batch → Task → Annotation",
            "REST APIs (OpenAPI): CRUD for all entities, queue, pipeline controls",
            "FIFO queue with deterministic claim locking (no double-claim)",
            "Multi-stage pipeline: L1 → Review → Done (configurable)",
            "Three portals: Ops (dashboard, controls), Rater (queue, annotate), Reviewer (review, approve)",
            "Auth & RBAC: Login, roles (Ops, Rater, Reviewer, Admin), internal vs external separation",
            "Task UI: Schema-driven; text, image, audio, video; single/multi-select, free-text",
            "Storage: PostgreSQL (transactional), S3/Azure Blob (assets), signed URLs",
            "Baseline logging: claims, submissions, status transitions (structured JSON)",
        ],
    )

    # ----- Slide 3: High-Level Architecture -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "High-Level Architecture (Phase 1)")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Client: Browser — Ops Portal, Rater Portal, Reviewer Portal (HTTPS)",
            "API layer: FastAPI — REST, OpenAPI, JWT auth, RBAC. Routes: /auth, /workspaces, /projects, /batches, /tasks, /queue",
            "Data layer:",
            "  • PostgreSQL — users, workspaces, projects, batches, tasks, annotations, pipeline state",
            "  • S3 / Azure Blob — task assets, uploads, export artifacts; signed URLs for access",
            "  • Redis (optional) — claim locking, cache, Celery broker for async jobs",
            "Deployment: Containerized API (Docker); cloud compute (Azure App Service / AWS ECS); secrets from vault",
        ],
    )

    # ----- Slide 4: Data Model -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Data Model — Core Entities")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Workspace  (1) ──────< (N) Project   — Top-level isolation; name, description",
            "Project    (1) ──────< (N) Batch    — pipeline_stages[], response_schema (JSON)",
            "Batch      (1) ──────< (N) Task     — FIFO order by created_at",
            "Task       (1) ──────< (N) Annotation — status, pipeline_stage, content (JSON), claimed_by_id",
            "User       (1) ──────< (N) Task (claimed)  — email, role, hashed_password",
            "User       (1) ──────< (N) Annotation     — response (JSON), pipeline_stage",
            "",
            "Key design: Response schema is JSON (flexible); UI components derived from schema.",
            "Pipeline stages and statuses are configurable per project (e.g. L1, Review, Hold, Done).",
        ],
        bullet=False,
    )

    # ----- Slide 5: Pipeline & Status Flow -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Pipeline & Status Flow")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Task statuses: pending → claimed → in_review → done (or hold, archived)",
            "Pipeline stages (e.g.): L1 → Review → Done. Ops can advance, hold, or archive.",
            "Flow: Rater claims task (FIFO) → submits annotation → task moves to Review.",
            "       Reviewer opens task → approves or edits → task moves to Done.",
            "Claim locking: Single source of truth (DB or Redis); no task claimed by two raters.",
            "All state transitions are logged for audit.",
        ],
    )

    # ----- Slide 6: Workflow & Data Visibility -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Workflow & Data Visibility (PRD-Aligned)")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Workflows: Configurable graphs of nodes (form, review, API, DB, manual) and edges; unique node ID and instance ID per run.",
            "Annotation flow: Upload file (image/doc) → form node (add attributes) → submit → review node → task enters work queue.",
            "Work queue: Manual/review nodes assign tasks to queues; users see only tasks in their queue (by role or hierarchy).",
            "Task tab: One view with filters and sorting (status, workflow, date, assignee); data scoped by visibility.",
            "Data visibility: Admin sees all workflows, tasks, and data; other users see only tasks assigned to them or in their queue.",
            "Child workflows: A node can start a sub-workflow; completion is sequential (child completes then parent continues).",
        ],
    )

    # ----- Slide 7: Infrastructure -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Infrastructure (Phase 1)")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    add_body(
        box2,
        [
            "Compute:",
            "  • API: Containerized (Docker); Azure App Service / Container Apps or AWS ECS",
            "  • Optional workers: Celery + Redis for async (export, bulk ops)",
            "Database:",
            "  • PostgreSQL 14+ (managed: Azure DB for PostgreSQL or AWS RDS)",
            "  • JSONB for response_schema, content, response",
            "Storage:",
            "  • S3 or Azure Blob: task assets, uploads, export artifacts",
            "  • Signed URLs for secure access; no public bucket exposure",
        ],
    )
    box3 = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(5.8))
    add_body(
        box3,
        [
            "Secrets: Azure Key Vault / AWS Secrets Manager (DB URL, API keys, JWT secret)",
            "DNS & TLS: Hostnames and HTTPS via cloud LB or App Gateway",
            "Environments: Dev, Staging, Prod with isolated DB and storage per env",
        ],
    )

    # ----- Slide 8: Deployment & CI/CD -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Deployment & CI/CD")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Version control: Git, GitHub (or customer’s repo); branch strategy (e.g. main + feature branches)",
            "CI: GitHub Actions (or customer’s pipeline) — on push/PR: build, run tests, lint",
            "Build: Docker image for API (and optional workers); multi-stage build for smaller image",
            "Deploy: Push image to registry (ACR/ECR); deploy to App Service / ECS with rolling update",
            "Database: Migrations (e.g. Alembic) run as part of release or separate step; backward-compatible",
            "Secrets: Injected via environment or cloud secret store; never in image or repo",
            "Rollback: Previous image tag; DB migrations are forward and backward compatible where possible",
        ],
    )

    # ----- Slide 9: Monitoring, Logging & Insights -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Monitoring, Logging & Insights")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Logging: Structured JSON logs (request id, user, action, duration, status); stdout → log aggregator",
            "Events logged: login, task claim, task submit, status change, Ops actions (advance, archive); no PII in logs",
            "Metrics: API latency (p50/p95), error rate, queue depth, tasks per status; exposed via Prometheus or cloud native",
            "Dashboards: Operational view (tasks by status, throughput, errors); Datadog / Azure Monitor / CloudWatch",
            "Alerts: Error rate spike, latency SLO breach, queue backlog; PagerDuty/email/Slack as needed",
            "Audit: Key actions (who did what, when) stored and searchable; retention per policy (e.g. 90 days)",
        ],
    )

    # ----- Slide 10: Security & Compliance -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Security & Compliance (Phase 1)")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Authentication: Login with credentials; JWT for API; optional Okta SSO in Phase 2",
            "Authorization: RBAC — Ops, Rater, Reviewer, Admin; internal vs external separation; enforced server-side",
            "Data in transit: HTTPS only; TLS 1.2+",
            "Data at rest: DB and blob storage encrypted (cloud default or customer-managed keys)",
            "Secrets: Stored in vault; rotated per policy; not in code or config in repo",
            "Audit: All sensitive actions logged; logs retained and accessible for compliance",
        ],
    )

    # ----- Slide 11: Why We Can Deliver -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Why We Can Build This — Confidence")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Proven stack: Python/FastAPI, PostgreSQL, Redis, S3/Blob — widely used and maintainable",
            "Clear scope: Phase 1 is bounded by PRD; we deliver the MVP flow and foundations first",
            "Working demo: V3 demo shows end-to-end flow (Ops, Rater, Reviewer); we extend and harden it",
            "Single codebase: One platform for all personas; no fragmented systems",
            "Phased delivery: You see working software in 8–10 weeks; we iterate with your feedback",
            "Discovery first: We align on users, load, and integrations before build to avoid rework",
            "Workflow & visibility: Graph-based workflows and Admin vs user-wise data are designed per PRD (see Developer Guide).",
        ],
    )

    # ----- Slide 12: Next Steps -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Next Steps")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "1. Align on Phase 1 scope and timeline (8–10 weeks, 3–4 people)",
            "2. Confirm infrastructure: cloud, region, and any existing CI/CD we plug into",
            "3. Answer discovery questions (users, load, integrations, data constraints)",
            "4. Kick off: repo, backlog, and first sprint; regular demos to stakeholders",
        ],
    )

    import os
    out_path = os.path.join(os.path.dirname(__file__) or ".", "Phase1_Data_Annotation_Platform.pptx")
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    main()
