"""
Generate FINAL Data Annotation Smart Factory Proposal PPT.
Based on teammate's proposal; enhanced with technical depth (tech stack, infra, architecture, database).
Timeline, resources, and profiles unchanged.

Run from this (demo) folder: python generate_final_proposal_ppt.py
Output: Data_Annotation_Smart_Factory_Proposal_Final.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE_COLOR = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_COLOR = RGBColor(0x16, 0xa3, 0x4a)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
BULLET_COLOR = RGBColor(0x55, 0x55, 0x55)


def add_body(shape, lines, bullet=True, font_size=14):
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BULLET_COLOR if bullet else BODY_COLOR
        p.space_after = Pt(4)
        if bullet:
            p.level = 0


def add_section_title(tf, text):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(10)


def slide_title(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = BODY_COLOR
        p2.space_before = Pt(8)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ----- Slide 1: Title (same as teammate) -----
    slide = slide_title(
        prs,
        "The Next-Gen Data Annotation Platform",
        "Phase 1 MVP Proposal: A Scalable, Serverless 'Smart Factory'",
    )
    box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Presented By: Anil Rawat  |  Date: February 6, 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = BULLET_COLOR

    # ----- Slide 2: Executive Summary -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Executive Summary — Replacing Manual Chaos with Digital Precision")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Goal: Transition from distributed, manual workflows to a centralized internal platform.",
            "Deliverable: Fully operational Phase 1 MVP in 8 Weeks.",
            "Strategy: 'Smart Factory' approach using Serverless Event-Driven Architecture.",
            "Key Win: A 'No-Ops' system that eliminates maintenance costs while handling high-volume concurrency.",
            "Technical Foundation: REST APIs (OpenAPI), FIFO queue with deterministic claim locking, hybrid DB (PostgreSQL + MongoDB).",
        ],
    )

    # ----- Slide 3: Understanding Your Challenge -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Understanding Your Challenge — Why the Current Process is Broken")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "Fragmented Workflows: Operations distributed across multiple tools create inefficiency.",
            "Manual Bottlenecks: Ops Managers act as 'traffic cops' rather than quality managers.",
            "Vendor Dependency: Reliance on external vendors for simple tasks due to lack of tooling.",
            "Risk: High-volume projects risk data corruption ('double-work') without deterministic locking.",
        ],
    )

    # ----- Slide 4: Our Vision -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Our Vision: The 'Zero-Touch' Pipeline")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "We are not just building a tool; we are building an autonomous engine.",
            "Current State: Manual file hand-offs, email coordination, rigid vendor tools.",
            "Future State: A 'Digital Assembly Line' where data flows automatically from Ingest -> Label -> Review.",
            "The Core Promise: Deterministic Locking. No data is ever lost, no task is ever doubled.",
            "The Automation Advantage: We use 'Robotic Gatekeepers' (Linters) to reject bad data instantly, reducing human review time by 40%.",
        ],
    )

    # ----- Slide 5: Solution — Smart Factory -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Solution: The 'Smart Factory' — An Automated Assembly Line for Data")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    add_body(
        box2,
        [
            "1. The Traffic Controller (Queueing): AWS SQS FIFO queues guarantee order; handle 10,000+ concurrent users; exactly-once processing with message deduplication.",
            "2. The Robotic Gatekeeper (Auto-QA): Serverless Linters (Lambda) validate every submission instantly. Bad data is rejected at the source; schema validation + business rules.",
            "3. Dual-Vault Storage: PostgreSQL (transactional state, claim locking, ACID) + MongoDB (flexible annotation payloads, schema evolution) + S3 (assets, signed URLs).",
        ],
    )

    # ----- Slide 6: Technology Stack (enhanced technical) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Technology Stack — Built on AWS (Serverless, No-Ops)")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(5.9))
    add_body(
        box2,
        [
            "COMPUTE: AWS Lambda (Node.js 20.x). Reserved concurrency per queue consumer. VPC for DB access. $0 when idle.",
            "QUEUEING: AWS SQS FIFO. Message group ID = batch; dedup window 5 min. 300 msg/s per queue; scale via multiple queues.",
            "STORAGE:",
            "  • PostgreSQL 15 (RDS or Aurora): users, workspaces, projects, batches, tasks, claim state; JSONB for config; row-level locking (SELECT FOR UPDATE SKIP LOCKED).",
            "  • MongoDB Atlas: annotation documents (flexible schema); TTL indexes for audit; read replicas for reporting.",
            "  • S3: task assets; bucket private; pre-signed URLs (15 min expiry); lifecycle to Glacier for archive.",
        ],
        font_size=12,
    )
    box3 = slide.shapes.add_textbox(Inches(5.2), Inches(1.15), Inches(4.3), Inches(5.9))
    add_body(
        box3,
        [
            "SECURITY: Okta (OIDC, JWT); AWS WAF (rate limit, geo, bot); Secrets Manager for DB credentials; no secrets in code.",
            "API: API Gateway REST (OpenAPI 3); request validation; CORS; throttling 10k req/s.",
            "FRONTEND: React 18 + Tailwind CSS; static build; CloudFront + S3; global CDN; HTTPS only.",
            "OBSERVABILITY: CloudWatch Logs (structured JSON); X-Ray tracing; alarms on Lambda errors and queue depth.",
        ],
        font_size=12,
    )

    # ----- Slide 7: Architecture Blueprint (technical) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Architecture Blueprint (Phase 1) — Technical Layers")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(5.9))
    add_body(
        box2,
        [
            "Layer 1 — Edge: CloudFront (HTTPS, TLS 1.2+); WAF; DDoS protection.",
            "Layer 2 — API: API Gateway REST; routes /auth, /workspaces, /projects, /batches, /tasks, /queue; JWT validation (Okta); internal vs external path separation.",
            "Layer 3 — Compute: Lambda (auth, queue consumer, linter, export); event source = SQS; async processing; DLQ for failed messages.",
            "Layer 4 — Data: RDS PostgreSQL (primary + standby); MongoDB Atlas (VPC peering or private endpoint); S3 (private buckets).",
            "Security: No public DB; Lambda in VPC; IAM least privilege; encryption at rest (KMS) and in transit.",
        ],
    )

    # ----- Slide 8: Database & Data Model (technical) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Database & Data Model — Concurrency & Integrity")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(5.9))
    add_body(
        box2,
        [
            "PostgreSQL (state of record):",
            "  • Tables: users, workspaces, projects, batches, tasks, annotations.",
            "  • tasks: id, batch_id, status, claimed_by_id, claimed_at, pipeline_stage, created_at.",
            "  • Claim lock: UPDATE tasks SET claimed_by_id = $1, claimed_at = NOW() WHERE id IN (SELECT id FROM tasks WHERE batch_id = $2 AND claimed_by_id IS NULL ORDER BY created_at LIMIT 1 FOR UPDATE SKIP LOCKED) RETURNING *.",
            "  • Indexes: (batch_id, status), (claimed_by_id), (created_at); JSONB GIN for filters.",
        ],
        font_size=11,
    )
    box3 = slide.shapes.add_textbox(Inches(5.2), Inches(1.15), Inches(4.3), Inches(5.9))
    add_body(
        box3,
        [
            "MongoDB (annotation payloads):",
            "  • Collections: annotations (task_id, user_id, response JSON, pipeline_stage, created_at).",
            "  • Flexible schema per project response_schema; versioned for audit.",
            "S3: Key = workspace/project/batch/asset_id; metadata for content-type; pre-signed GET/PUT.",
            "Atomicity: Task claim + annotation write in same transaction where possible; else idempotent keys.",
        ],
        font_size=11,
    )

    # ----- Slide 9: Infrastructure & DevOps -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Infrastructure & DevOps — CI/CD, IaC, Environments")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(5.9))
    add_body(
        box2,
        [
            "IaC: AWS CDK or Terraform; VPC, Lambda, SQS, RDS, S3, API Gateway, WAF; versioned in Git.",
            "CI/CD: GitHub Actions (or AWS CodePipeline); on push to main: lint, test, build; deploy to Staging; manual promote to Prod.",
            "Environments: Dev (developer laptops + shared Dev AWS); Staging (mirrors Prod, for UAT); Prod (isolated, no direct DB access).",
            "Secrets: AWS Secrets Manager; rotated for DB; Lambda fetches at cold start; no .env in repo.",
            "Monitoring: CloudWatch dashboards (queue depth, Lambda duration, error rate); alarms → SNS → email/Slack; log retention 90 days.",
        ],
    )

    # ----- Slide 10: User Experience (same two portals) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "The User Experience: Two Distinct Portals")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(5.9))
    add_body(
        box2,
        [
            "1. The Ops Command Center (Internal Manager)",
            "  • Dashboard with real-time 'Traffic Light' status (Green=On Track, Red=Blocked).",
            "  • Drag-and-drop batch upload; 'Emergency Brake' (Hold/Pause); One-Click Export.",
            "  • Total control without engineering support.",
            "2. The Rater Workbench (External Worker)",
            "  • Distraction-free UI; large central canvas; minimal toolbars.",
            "  • 'Next Task' auto-fetch; keyboard shortcuts; real-time Linter feedback.",
            "  • Optimized for throughput and ergonomic comfort.",
        ],
        font_size=12,
    )

    # ----- Slide 11: End-to-End Workflow -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "How It Works: End-to-End Workflow")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(5.9))
    add_body(
        box2,
        [
            "STEP 1 — INGEST: Ops creates project, uploads 10k images; system validates and pushes tickets to SQS Pending Queue.",
            "STEP 2 — ASSIGNMENT: Rater clicks 'Start'; API pops top task from Queue; PostgreSQL locks task to User A (SELECT FOR UPDATE SKIP LOCKED); no double-claim.",
            "STEP 3 — EXECUTION: Rater labels; auto-save (localStorage); on Submit, Linter Lambda validates; if invalid, submission blocked.",
            "STEP 4 — DELIVERY: Valid tasks move to Done/Review; Manager exports clean JSON/CSV with signed S3 links.",
        ],
    )

    # ----- Slide 12: Why Trust -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Why Trust This Solution?")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(5.9))
    add_body(
        box2,
        [
            "FINANCIAL SAFETY: Serverless = Zero Idle Costs. If the project pauses, AWS bill ~$0. OPEX, not CAPEX.",
            "DATA INTEGRITY: Atomic locking in PostgreSQL (FOR UPDATE SKIP LOCKED). Physically impossible for two raters to claim the same task.",
            "SCALABILITY: SQS buffers traffic. Same architecture for 5 or 5,000 users; no re-platforming.",
        ],
    )

    # ----- Slide 13: Team & Timeline (EXACT same as teammate) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "The 'Lean Elite' Team & Timeline")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.6), Inches(5.9))
    add_body(
        box2,
        [
            "The Team (Total: 4 Resources)",
            "  • 1x Solution Architect (Serverless Design & Security)",
            "  • 2x Senior Full-Stack Engineers (Frontend UI + Backend Queue)",
            "  • 1x UI/UX Designer (Design UI)",
        ],
    )
    box3 = slide.shapes.add_textbox(Inches(5.2), Inches(1.15), Inches(4.3), Inches(5.9))
    add_body(
        box3,
        [
            "The 8-Week Plan",
            "  • Week 1-2: Infrastructure Setup & Auth (Okta)",
            "  • Week 3-6: 'Traffic Controller' Queue & Task UI",
            "  • Week 7: End-to-End Testing (5k user sim) & Ops Tools",
            "  • Week 8: UAT & Handoff",
        ],
    )

    # ----- Slide 14: Pricing & Efficiency (same) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Pricing & Efficiency Model")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(5.9))
    add_body(
        box2,
        [
            "Development Cost: Optimized via AI-Augmented coding (30% faster cycles).",
            "Operational Cost (OPEX):",
            "  • Traditional Server Approach: High fixed monthly cost ($500-$1000+).",
            "  • Our Serverless Approach: Pay-per-ms (~$50-$100/mo depending on traffic).",
            "ROI: System pays for itself via reduced DevOps hours and eliminated idle server costs.",
        ],
    )

    # ----- Slide 15: Next Steps (same) -----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_section_title(box.text_frame, "Next Steps — Ready to Launch Your Factory?")
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(5.9))
    add_body(
        box2,
        [
            "1. Immediate Action: Approval of SOW for Phase 1 MVP.",
            "2. Pre-Requisites: Access to AWS Staging Environment & Okta Config.",
            "3. Commitment: Functional, secure platform deployed in 45 days.",
        ],
    )

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)) or ".", "Data_Annotation_Smart_Factory_Proposal_Final.pptx")
    prs.save(out_path)
    print("Generated:", out_path)
    return out_path


if __name__ == "__main__":
    main()
